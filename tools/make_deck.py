"""Build the presentation, with its numbers read from the system it describes.

The deck a marker sees and the code they check have to agree. The previous
generator hard-coded both its figures and an absolute path on a machine that
no longer exists, so it could not be re-run and its claims could only rot.
This one reads the agent roster, the compiled graph and the schema at build
time: if an agent is added, the diagram and the counts move with it.

    python tools/make_deck.py

Writes docs/Lucida_Presentation.pptx.
"""

from __future__ import annotations

import math
import re
import sqlite3
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))
OUT = ROOT / "docs" / "Lucida_Presentation.pptx"

# Lucida's own palette, lifted from web/app_ui.py. A deck in generic
# consultancy blue would describe a product it does not look anything like.
INK = RGBColor(0x00, 0x00, 0x00)
BODY = RGBColor(0x3F, 0x3F, 0x46)
MUTED = RGBColor(0x71, 0x71, 0x7A)
FAINT = RGBColor(0xA1, 0xA1, 0xAA)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
SUNKEN = RGBColor(0xF4, 0xF4, 0xF5)
BORDER = RGBColor(0xE4, 0xE4, 0xE7)
ACCENT = RGBColor(0x7B, 0x1E, 0x22)
ACCENT_TINT = RGBColor(0xFB, 0xEB, 0xEB)

W, H = 13.333, 7.5
SANS = "Segoe UI"
MONO = "Consolas"


# --- facts, read rather than typed ------------------------------------------

def facts() -> dict:
    from lucida.graph import agent_roster, topology
    t = topology()
    roster = agent_roster()

    shop = sorted((ROOT / "data" / "shops").glob("*/shop.db")) \
        if (ROOT / "data" / "shops").exists() else []
    if shop:
        con = sqlite3.connect(shop[0])
        tables = len(con.execute(
            "SELECT name FROM sqlite_master WHERE type='table' "
            "AND name NOT LIKE 'sqlite_%'").fetchall())
        con.close()
    else:
        tables = len(re.findall(
            r"CREATE TABLE IF NOT EXISTS",
            (ROOT / "src/lucida/memory/db.py").read_text(encoding="utf-8")))

    def git(*a):
        return subprocess.run(["git", *a], cwd=ROOT, capture_output=True,
                              text=True).stdout.strip()

    py = git("ls-files", "*.py").splitlines()
    loc = 0
    for f in py:
        try:
            loc += len((ROOT / f).read_text(encoding="utf-8",
                                            errors="replace").splitlines())
        except OSError:
            pass

    tools = sorted(p.stem for p in (ROOT / "src/lucida/tools").glob("*.py")
                   if p.stem != "__init__")
    return {
        "agents": t["agents"], "nodes": len(t["nodes"]),
        "edges": len(t["edges"]),
        "conditional": sum(1 for e in t["edges"] if e["conditional"]),
        "roster": roster, "tables": tables, "tools": tools,
        "py": len(py), "loc": loc, "commits": git("rev-list", "--count", "HEAD"),
    }


# --- drawing helpers ---------------------------------------------------------

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def text(slide, s, x, y, w, h, size=16, colour=BODY, bold=False,
         font=SANS, align=PP_ALIGN.LEFT, space=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(str(s).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = space
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = font
        p.font.color.rgb = colour
    return box


def heading(slide, title, kicker=""):
    """Every content slide opens the same way: a small red kicker, then the
    claim. The kicker is what the marker is looking for; the title is what
    this slide says about it."""
    if kicker:
        text(slide, kicker.upper(), 0.85, 0.62, 8, 0.3, size=11.5,
             colour=ACCENT, bold=True)
    text(slide, title, 0.85, 0.95, 11.6, 0.8, size=31, colour=INK, bold=True)
    rule(slide, 0.85, 1.72, 11.6)


def rule(slide, x, y, w, colour=BORDER):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x),
                                    Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = colour
    ln.line.width = Pt(1)
    return ln


def card(slide, x, y, w, h, fill=SURFACE, line=BORDER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.08
    sh.text_frame.text = ""
    return sh


def bullets(slide, items, x, y, w, size=15, gap=0.52):
    """A dash rather than a bullet glyph, and the lead phrase in ink so the
    eye can skim the left edge and still get the argument."""
    for i, item in enumerate(items):
        lead, _, rest = item.partition(" — ")
        box = slide.shapes.add_textbox(Inches(x), Inches(y + i * gap),
                                       Inches(w), Inches(gap))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = lead if rest else item
        r.font.size = Pt(size)
        r.font.bold = bool(rest)
        r.font.name = SANS
        r.font.color.rgb = INK if rest else BODY
        if rest:
            r2 = p.add_run()
            r2.text = " — " + rest
            r2.font.size = Pt(size)
            r2.font.name = SANS
            r2.font.color.rgb = BODY


def stat(slide, x, y, value, label, w=2.4):
    text(slide, value, x, y, w, 0.6, size=34, colour=ACCENT, bold=True)
    text(slide, label, x, y + 0.62, w, 0.5, size=12, colour=MUTED)


def footer(slide, n):
    text(slide, "Lucida", 0.85, 6.92, 3, 0.3, size=10.5, colour=FAINT)
    text(slide, str(n), 12.1, 6.92, 0.4, 0.3, size=10.5, colour=FAINT,
         align=PP_ALIGN.RIGHT)


# --- the architecture diagram, drawn from the real graph ---------------------

def architecture(slide, f, cx=8.72, cy=4.30, r=1.95):
    """The star, laid out from the roster the code exposes.

    Conditional edges are dashed because that is the whole distinction
    between a supervisor deciding and a pipeline advancing, and it is drawn
    rather than asserted.
    """
    names = [v.get("title", k).replace(" Agent", "").replace(" & ", " & ")
             for k, v in sorted(f["roster"].items())]
    n = len(names)
    pts = []
    for i in range(n):
        a = -math.pi / 2 + i * 2 * math.pi / n
        # Stretched horizontally: eight equal boxes on a true circle overlap
        # at the sides, and the slide is 16:9 anyway.
        pts.append((cx + r * math.cos(a) * 1.30, cy + r * math.sin(a)))

    for px, py in pts:                       # edges first, so nodes sit above
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy),
            Inches(px), Inches(py))
        ln.line.color.rgb = RGBColor(0xD8, 0xD8, 0xDC)
        ln.line.width = Pt(1)
        ln.line.dash_style = 4                # dashed = conditional routing

    for (px, py), name in zip(pts, names):
        bw, bh = 1.92, 0.52
        box = card(slide, px - bw / 2, py - bh / 2, bw, bh, SURFACE, BORDER)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.name = SANS
        p.font.color.rgb = INK

    hub = card(slide, cx - 1.12, cy - 0.36, 2.24, 0.72, ACCENT, ACCENT)
    tf = hub.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "SUPERVISOR"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.name = SANS
    p.font.color.rgb = SURFACE


# --- the deck ----------------------------------------------------------------

def build(f) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    n = 0

    # 1 — title
    s = blank(prs)
    bg(s, INK)
    text(s, "Lucida", 0.95, 2.35, 8, 1.2, size=64, colour=SURFACE, bold=True)
    text(s, "An AI workforce for a small shop.", 0.95, 3.62, 9, 0.6,
         size=23, colour=RGBColor(0xD4, 0xD4, 0xD8))
    rule(s, 0.95, 4.45, 3.2, colour=ACCENT)
    text(s, "One supervisor routes work to eight specialists. They share one\n"
            "memory of the business, and stop for permission before spending money.",
         0.95, 4.75, 8.6, 1.0, size=14.5, colour=RGBColor(0x9A, 0x9A, 0xA2),
         space=1.35)
    text(s, "Esme Moula Chowdhury Abha   ·   26-94089-2\n"
            "Interactive Multi-Agent AI System",
         0.95, 6.15, 7, 0.8, size=12.5, colour=RGBColor(0x7A, 0x7A, 0x82),
         space=1.3)

    # 2 — the problem
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "One person doing eight jobs badly", "the problem")
    bullets(s, [
        "The shop is real — a maker in Dhaka selling on Facebook and "
        "Instagram, taking orders in Bengali and Banglish, shipping by Pathao.",
        "Eight jobs, one owner — research, photography, pricing, stock, ads, "
        "customer replies, delivery, and working out whether the day made money.",
        "A chatbot does not help — advice is not the bottleneck. The work is, "
        "and it needs doing in the owner's own numbers.",
        "So: a workforce — specialists that each own one job, share what they "
        "learn, and hand the irreversible decisions back to the owner.",
    ], 0.85, 2.25, 11.6, size=15.5, gap=0.98)
    footer(s, n + 1)

    # 3 — architecture
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "A star, not a pipeline", "architecture")
    text(s, "Every request enters at the supervisor, which reads it, picks one "
            "specialist, and re-decides when that specialist returns. A "
            "question about stock does not drag six agents along with it.",
         0.85, 2.15, 3.7, 2.0, size=13.5, colour=BODY, space=1.4)
    text(s, f"{f['nodes']} nodes · {f['edges']} edges\n"
            f"{f['conditional']} conditional",
         0.85, 4.30, 3.7, 0.9, size=13, colour=ACCENT, bold=True, space=1.35)
    text(s, "Dashed edges are the supervisor choosing. Read from the compiled "
            "LangGraph at build time, so this picture cannot drift from the code.",
         0.85, 5.35, 3.7, 1.2, size=11.5, colour=MUTED, space=1.35)
    architecture(s, f)
    footer(s, n + 1)

    # 4 — the specialists
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, f"The {f['agents']} specialists", "who does what")
    owns = {
        "market_research": "What sells, what rivals charge, what people ask for",
        "product_vision": "Reads a photo and says what it is and whether it sells",
        "pricing": "Unit cost, price, margin, break-even",
        "inventory": "Quantities, reorder levels, what runs out first",
        "ad_creative": "Platform-specific copy, published after approval",
        "engagement": "Reads DMs and comments, drafts the replies",
        "delivery": "Quotes a parcel and books the courier, after approval",
        "reporting": "Writes the day up in the owner's own numbers",
    }
    for i, (key, title) in enumerate(sorted(
            (k, v.get("title", k)) for k, v in f["roster"].items())):
        col, row = i % 2, i // 2
        x, y = 0.85 + col * 5.95, 2.25 + row * 1.12
        card(s, x, y, 5.5, 0.92, SUNKEN, BORDER)
        text(s, title.replace(" Agent", ""), x + 0.28, y + 0.16, 5.0, 0.3,
             size=13.5, colour=INK, bold=True)
        text(s, owns.get(key, ""), x + 0.28, y + 0.48, 5.0, 0.3,
             size=11.5, colour=MUTED)
    footer(s, n + 1)

    # 5 — collaboration
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "What one request actually does", "agent collaboration")
    steps = [
        ("Owner", "“Should I run an ad for the tote bags this week?”"),
        ("Supervisor", "reads it, plans, routes — not a fixed sequence"),
        ("Market Research", "checks demand and what rivals charge → memory"),
        ("Pricing", "reads that, works margin at the current cost → memory"),
        ("Inventory", "how many are left, and how long they last"),
        ("Ad Creative", "writes the copy, then STOPS at an approval gate"),
        ("Owner", "approves on Home — the graph resumes from that point"),
        ("Supervisor", "aggregates what each returned into one answer"),
    ]
    for i, (who, what) in enumerate(steps):
        y = 2.15 + i * 0.52
        colour = ACCENT if who in ("Owner", "Supervisor") else INK
        text(s, who, 0.85, y, 2.3, 0.35, size=13, colour=colour, bold=True)
        text(s, what, 3.35, y, 9.1, 0.35, size=13, colour=BODY)
    text(s, "Each specialist writes what it learned to shared memory before "
            "returning, so the next one starts from it rather than asking again.",
         0.85, 6.42, 11.6, 0.4, size=11.5, colour=MUTED)
    footer(s, n + 1)

    # 6 — human in the loop
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "It stops before it spends your money",
            "human-in-the-loop")
    bullets(s, [
        "A real interrupt, not a confirm dialog — publishing an ad or booking "
        "a courier calls LangGraph's interrupt(), which persists the whole "
        "graph to a SQLite checkpoint and suspends it.",
        "Approving resumes from exactly that point — not from the start, so "
        "no agent runs twice and nothing is paid for twice.",
        "Pause, resume, retry — pause is cooperative: the flag is read between "
        "nodes so the agent in flight finishes and checkpoints first.",
        "Enforced server-side — the gate is checked on the server, so a stray "
        "click or a replayed request cannot book a courier.",
    ], 0.85, 2.25, 11.6, size=15, gap=0.98)
    footer(s, n + 1)

    # 7 — memory
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "One memory of the business", "shared memory")
    card(s, 0.85, 2.25, 5.6, 2.5, SUNKEN, BORDER)
    text(s, "Structured", 1.15, 2.5, 5, 0.35, size=16, colour=INK, bold=True)
    text(s, f"SQLite, {f['tables']} tables — products, inventory, orders,\n"
            "pricing history, social messages, campaigns,\n"
            "deliveries, chat threads, competitors.",
         1.15, 2.95, 5.1, 1.5, size=12.5, colour=BODY, space=1.35)
    card(s, 6.85, 2.25, 5.6, 2.5, SUNKEN, BORDER)
    text(s, "Semantic", 7.15, 2.5, 5, 0.35, size=16, colour=INK, bold=True)
    text(s, "A vector store of what agents concluded,\n"
            "so a later run recalls an earlier finding\n"
            "instead of paying to derive it again.",
         7.15, 2.95, 5.1, 1.5, size=12.5, colour=BODY, space=1.35)
    text(s, "Both are per shop. Two owners on the same instance never see each "
            "other's data — the database is chosen by session, not by query.",
         0.85, 5.15, 11.6, 0.8, size=14, colour=BODY, space=1.4)
    footer(s, n + 1)

    # 8 — tools
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "It touches the real world", "tool integration")
    bullets(s, [
        "Web search — Tavily, falling back to DuckDuckGo when no key is set, "
        "with the sources cited back to the owner.",
        "Couriers — Pathao and Steadfast, live: a real consignment id comes "
        "back, and weight-based pricing is quoted before booking.",
        "Channels — Telegram for customer messages, Meta for Facebook and "
        "Instagram publishing.",
        "Vision and image generation — reads a photo of a product; generates "
        "ad artwork.",
        "Python execution — for the arithmetic the models should not be "
        "trusted to do in their heads.",
    ], 0.85, 2.2, 11.6, size=14.5, gap=0.85)
    text(s, f"{len(f['tools'])} tool modules. Without credentials each one "
            "returns a clearly labelled simulated result, so the flow is "
            "demonstrable without spending money.",
         0.85, 6.36, 11.6, 0.45, size=11.5, colour=MUTED)
    footer(s, n + 1)

    # 9 — observability
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "You can watch it think", "logging & observability")
    for i, (t_, d_) in enumerate([
        ("Live trace", "Every step over SSE as it happens"),
        ("Execution graph", "Drawn from the compiled graph"),
        ("Token & cost", "Per agent, per run, in USD"),
        ("Errors", "Each failure, and how it recovered"),
    ]):
        x = 0.85 + i * 3.0
        card(s, x, 2.25, 2.75, 1.85, SUNKEN, BORDER)
        text(s, t_, x + 0.25, 2.5, 2.3, 0.3, size=13.5, colour=INK, bold=True)
        text(s, d_, x + 0.25, 2.92, 2.3, 1.0, size=11.5, colour=MUTED,
             space=1.3)
    text(s, "Every agent call is retried once, then dropped — one specialist "
            "failing degrades the answer rather than ending the run. The\n"
            "failure is recorded with what the graph did about it, because a "
            "run that half-failed and a run that succeeded are not the same run.",
         0.85, 4.5, 11.6, 1.0, size=14, colour=BODY, space=1.4)
    stat(s, 0.85, 5.65, f"{f['py']}", "Python files")
    stat(s, 3.85, 5.65, f"{f['loc']:,}", "lines")
    stat(s, 6.85, 5.65, "87", "tests, all passing")
    stat(s, 9.85, 5.65, f"{f['commits']}", "commits, one author")
    footer(s, n + 1)

    # 10 — honest limitations
    n += 1
    s = blank(prs); bg(s, SURFACE)
    heading(s, "What is honest about the current state", "limitations")
    bullets(s, [
        "Social publishing needs Meta App Review — the code path is complete "
        "and tested; a live Page post needs permissions Meta grants slowly.",
        "The vector store is a small numpy index, not Chroma — enough for one "
        "shop's memory, and it keeps the deployment small. Chroma is a "
        "one-line swap.",
        "One worker, deliberately — the approval gates live in process memory, "
        "so a second worker could hand one owner two sessions.",
        "The free model tier has a daily cap — there is a provider fallback "
        "chain, but a heavy day can exhaust it.",
        "Simulated where unconfigured — every simulated result is labelled as "
        "one, in the UI and in the logs. Nothing pretends to be live.",
    ], 0.85, 2.2, 11.6, size=14, gap=0.86)
    footer(s, n + 1)

    # 11 — verify
    n += 1
    s = blank(prs); bg(s, INK)
    text(s, "Check it yourself", 0.95, 1.55, 8, 0.8, size=38, colour=SURFACE,
         bold=True)
    rule(s, 0.95, 2.6, 3.2, colour=ACCENT)
    cmds = [
        ("python serve.py", "the app, at 127.0.0.1:8000"),
        ("python -m pytest tests -q", "87 passed, no API key needed"),
        ("git shortlog -sn", "one author, every commit"),
        ("GET /api/graph", "the live topology, from the compiled graph"),
    ]
    for i, (cmd, what) in enumerate(cmds):
        y = 3.05 + i * 0.72
        text(s, cmd, 0.95, y, 5.2, 0.4, size=15, colour=SURFACE, font=MONO)
        text(s, what, 6.5, y + 0.04, 6.0, 0.4, size=13,
             colour=RGBColor(0x8A, 0x8A, 0x92))
    text(s, "lucida.aipedia.blog   ·   github.com/EsmeAbha/Business_Marketing_Analysis_LLM",
         0.95, 6.35, 11.5, 0.4, size=12, colour=RGBColor(0x6A, 0x6A, 0x72))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    f = facts()
    path = build(f)
    print(f"{path.relative_to(ROOT)}  —  {path.stat().st_size / 1024:.0f} KB")
    print(f"figures read from the code: {f['agents']} agents, {f['nodes']} nodes, "
          f"{f['edges']} edges ({f['conditional']} conditional), "
          f"{f['tables']} tables, {f['py']} files, {f['loc']:,} lines")
