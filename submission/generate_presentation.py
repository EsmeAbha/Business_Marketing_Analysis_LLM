from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

root = Path(r'd:\business-analysis-LLM\Business_Marketing_Analysis_LLM')
assets = root / 'submission' / 'slide_assets'
output = root / 'submission' / 'Lucida_Presentation.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

navy = RGBColor(15, 23, 42)
teal = RGBColor(20, 184, 166)
blue = RGBColor(37, 99, 235)
gray = RGBColor(107, 114, 128)
light = RGBColor(243, 244, 246)
white = RGBColor(255, 255, 255)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, x=0.6, y=0.45, w=8.5, h=0.7, color=white):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, navy)
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
bar.fill.solid(); bar.fill.fore_color.rgb = teal
bar.line.fill.background()

logo = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(5.0), Inches(0.8))
logo_tf = logo.text_frame
p = logo_tf.paragraphs[0]
p.text = 'Lucida'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = white

sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(9.0), Inches(0.8))
sub_tf = sub.text_frame
p = sub_tf.paragraphs[0]
p.text = 'Interactive Multi-Agent AI Workforce for a Small Business'
p.font.size = Pt(25)
p.font.color.rgb = white

bullets = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.5), Inches(2.8))
btf = bullets.text_frame
items = [
    'Supervisor-driven business workflow',
    '8 specialist agents',
    'Shared memory + tool integration',
    'Human approval for risky actions',
    'Live dashboard + execution trace' 
]
for i, item in enumerate(items):
    p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
    p.text = '• ' + item
    p.font.size = Pt(20)
    p.font.color.rgb = white
    p.level = 0

img = assets / 'image1.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(7.2), Inches(1.9), Inches(5.2), Inches(3.8))

foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(8.0), Inches(0.4))
ft = foot.text_frame
p = ft.paragraphs[0]
p.text = 'AI Workforce | Business Intelligence | LangGraph'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(191, 219, 254)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, light)
add_title(slide, 'Problem and Motivation', 0.7, 0.5, 8, 0.7, navy)

left = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.7), Inches(4.0))
ltf = left.text_frame
items = [
    'Small businesses need support for research, pricing, stock, marketing, customer communication, and delivery.',
    'Single-chatbot tools cannot coordinate end-to-end decisions across multiple tasks.',
    'A real business system must reason, act, and wait for approval before spending money or publishing content.',
    'The system must be explainable, traceable, and grounded in shared memory.'
]
for i, item in enumerate(items):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    p.text = item
    p.font.size = Pt(21)
    p.font.color.rgb = navy
    p.space_after = Pt(10)

box = slide.shapes.add_shape(1, Inches(7.0), Inches(2.0), Inches(5.3), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(224, 242, 254)
box.line.color.rgb = blue
box.line.width = Pt(1.2)
box_tf = box.text_frame
box_tf.word_wrap = True
box_tf.margin_left = Inches(0.3)
box_tf.margin_right = Inches(0.3)
for text, size, bold in [
    ('Business workflow automation', 24, True),
    ('Research → pricing → stock → marketing → delivery', 18, False),
    ('One supervisor, many specialists', 18, False),
]:
    p = box_tf.paragraphs[0] if not box_tf.paragraphs[0].text else box_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = navy
    p.space_after = Pt(10)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, navy)
add_title(slide, 'System Architecture', 0.6, 0.5, 8.0, 0.7, white)
shape = slide.shapes.add_shape(1, Inches(0.9), Inches(1.6), Inches(11.6), Inches(4.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(17, 24, 39)
shape.line.color.rgb = teal; shape.line.width = Pt(1.2)
shape.text = 'Owner Request\n\nSupervisor Agent\n\nMarket Research | Product Vision | Pricing | Inventory\n\nAd Creative | Customer Engagement | Delivery | Reporting\n\nShared Memory + Vector Store + SQLite + Logging\n\nFinal Business Recommendation'
shape.text_frame.paragraphs[0].font.size = Pt(19)
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.color.rgb = white
for p in shape.text_frame.paragraphs[1:]:
    p.font.size = Pt(15)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, light)
add_title(slide, 'Specialized Agents', 0.7, 0.5, 8.0, 0.7, navy)

cards = [
    ('Market Research', 'Demand, competition, pricing gaps'),
    ('Product Vision', 'Image analysis and product fit'),
    ('Pricing', 'Cost, margin, break-even'),
    ('Inventory', 'Stock tracking and reorder alerts'),
    ('Ad Creative', 'Campaign copy and ad generation'),
    ('Customer Engagement', 'Message analysis and sentiment'),
    ('Delivery', 'Courier quotes and bookings'),
    ('Reporting', 'Summaries and recommendations'),
]
for idx, (title, body) in enumerate(cards):
    x = Inches(0.7 + (idx % 4) * 3.1)
    y = Inches(1.6 + (idx // 4) * 2.15)
    sh = slide.shapes.add_shape(1, x, y, Inches(2.7), Inches(1.55))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(239, 246, 255)
    sh.line.color.rgb = blue; sh.line.width = Pt(1.2)
    tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = navy
    p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(10); p2.font.color.rgb = gray

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, navy)
add_title(slide, 'Workflow and Decision Flow', 0.6, 0.5, 8.0, 0.7, white)
flow = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(4.2))
ftf = flow.text_frame
steps = [
    '1. User sends a request or uploads a product image',
    '2. Supervisor interprets intent and chooses the next agent',
    '3. Specialist performs the task and writes the result to shared memory',
    '4. Supervisor revisits routing until the goal is complete',
    '5. Risky operations pause for human approval',
    '6. Final answer is generated with evidence and numbers'
]
for i, step in enumerate(steps):
    p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
    p.text = step
    p.font.size = Pt(19)
    p.font.color.rgb = white
    p.space_after = Pt(10)

img2 = assets / 'image2.png'
if img2.exists():
    slide.shapes.add_picture(str(img2), Inches(9.2), Inches(1.8), Inches(3.3), Inches(2.5))

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, light)
add_title(slide, 'Interactive Dashboard and UI', 0.7, 0.5, 8.0, 0.7, navy)
left = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.7), Inches(4.2))
ltf = left.text_frame
items = [
    'Dashboard of active agents and workflow',
    'Live execution trace and event feed',
    'Agent communication history',
    'Execution graph visualisation',
    'Token usage and API cost estimation',
    'Logs, errors, and memory viewer',
    'Human-in-the-loop controls: pause, resume, approve, retry',
    'Final report and output viewer'
]
for i, item in enumerate(items):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    p.text = '• ' + item
    p.font.size = Pt(19)
    p.font.color.rgb = navy
    p.space_after = Pt(8)

img3 = assets / 'image3.png'
if img3.exists():
    slide.shapes.add_picture(str(img3), Inches(7.2), Inches(1.7), Inches(5.3), Inches(3.5))

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, navy)
add_title(slide, 'Tools, Memory, and Reliability', 0.7, 0.5, 8.5, 0.7, white)
left = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(3.8))
ltf = left.text_frame
items = [
    'Web search and local fallback search',
    'External business APIs',
    'Python code execution sandbox',
    'RAG and vector memory retrieval',
    'SQLite persistence and graph state',
    'Logging and error containment',
    'Owner approval gates before irreversible actions'
]
for i, item in enumerate(items):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    p.text = '• ' + item
    p.font.size = Pt(19)
    p.font.color.rgb = white
    p.space_after = Pt(9)

info = slide.shapes.add_shape(1, Inches(7.3), Inches(2.1), Inches(4.8), Inches(2.7))
info.fill.solid(); info.fill.fore_color.rgb = RGBColor(15, 118, 110)
info.line.color.rgb = white; info.line.width = Pt(1.2)
infotf = info.text_frame; infotf.word_wrap = True; infotf.margin_left = Inches(0.3); infotf.margin_right = Inches(0.3)
for text, size in [('Shared Memory', 22), ('Structured data + semantic retrieval', 17), ('Reliable coordination across agents', 17)]:
    p = infotf.paragraphs[0] if not infotf.paragraphs[0].text else infotf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, light)
add_title(slide, 'Conclusion', 0.8, 0.6, 8.0, 0.7, navy)
box = slide.shapes.add_shape(1, Inches(1.0), Inches(1.5), Inches(11.1), Inches(4.1))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(255, 255, 255)
box.line.color.rgb = RGBColor(148, 163, 184); box.line.width = Pt(1.2)
box_tf = box.text_frame; box_tf.word_wrap = True
texts = [
    ('Lucida demonstrates a real multi-agent AI system, not a simple chatbot.', 24, True),
    ('A Supervisor agent coordinates specialists, shares memory, and routes work based on the actual business task.', 20, False),
    ('The system is interactive, explainable, and designed for business decision-making with human oversight.', 20, False),
    ('Result: a practical AI workforce for small-business operations and workflow automation.', 20, True),
]
for text, size, bold in texts:
    p = box_tf.paragraphs[0] if not box_tf.paragraphs[0].text else box_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = navy
    p.space_after = Pt(12)

prs.save(output)
print(f'Created {output}')
